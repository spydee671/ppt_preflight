#!/usr/bin/env python3
"""
Generate a test .pptx that exercises all av_preflight checks.

Creates 5 slides covering: custom fonts, embedded image, linked video,
linked audio, transitions, animations, hyperlinks, speaker notes, hidden slide.

Usage: python make_test_pptx.py  →  writes test_deck.pptx
"""

import io
import struct
import zlib
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

OUT = 'test_deck.pptx'

VIDEO_RELTYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/video'
AUDIO_RELTYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio'

P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'


# ── Image helper ──────────────────────────────────────────────────────────────

def make_png(width=80, height=60, rgb=(70, 130, 180)):
    """Return bytes of a solid-colour PNG (no Pillow needed)."""
    r, g, b = rgb

    def chunk(ctype, data):
        body = ctype + data
        return struct.pack('>I', len(data)) + body + struct.pack('>I', zlib.crc32(body) & 0xffffffff)

    ihdr = struct.pack('>IIBBBBB', width, height, 8, 2, 0, 0, 0)  # 8-bit RGB
    row  = b'\x00' + bytes([r, g, b]) * width                      # filter=0, then pixels
    idat = zlib.compress(row * height)

    return b'\x89PNG\r\n\x1a\n' + chunk(b'IHDR', ihdr) + chunk(b'IDAT', idat) + chunk(b'IEND', b'')


# ── XML helpers ───────────────────────────────────────────────────────────────

def add_transition(slide, kind='fade', dur_ms=700, auto_advance_ms=None):
    """Inject a slide transition element."""
    sld = slide.element
    for old in sld.findall(qn('p:transition')):
        sld.remove(old)

    attribs = f'dur="{dur_ms}"'
    if auto_advance_ms is not None:
        attribs += f' advClick="0" advTm="{auto_advance_ms}"'

    xml = (f'<p:transition xmlns:p="{P_NS}" {attribs}>'
           f'<p:{kind}/></p:transition>')
    elem = etree.fromstring(xml)

    # Insert before p:timing or p:extLst (whichever comes first)
    insert_at = len(sld)
    for i, child in enumerate(sld):
        tag = child.tag.split('}')[-1]
        if tag in ('timing', 'extLst'):
            insert_at = i
            break
    sld.insert(insert_at, elem)


def add_animation(slide, shape_id):
    """Inject a minimal click-triggered fade-in animation on shape_id."""
    xml = f'''<p:timing xmlns:p="{P_NS}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst>
                      <p:cond evt="onBegin" delay="indefinite"/>
                    </p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" presetID="1" presetClass="entr" presetSubtype="0"
                               fill="hold" grpId="0" nodeType="clickEffect">
                          <p:stCondLst>
                            <p:cond delay="0"/>
                          </p:stCondLst>
                          <p:childTnLst>
                            <p:animEffect transition="in" filter="fade">
                              <p:cBhvr>
                                <p:cTn id="5" dur="500"/>
                                <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                              </p:cBhvr>
                            </p:animEffect>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst><p:cond evt="onPrevNode" delay="0"/></p:prevCondLst>
            <p:nextCondLst><p:cond evt="onNextNode" delay="0"/></p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>'''

    sld = slide.element
    for old in sld.findall(qn('p:timing')):
        sld.remove(old)
    elem = etree.fromstring(xml)
    insert_at = len(sld)
    for i, child in enumerate(sld):
        if child.tag.split('}')[-1] == 'extLst':
            insert_at = i
            break
    sld.insert(insert_at, elem)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide1_title(prs):
    """Title slide — Georgia + Calibri, speaker notes."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    title = slide.shapes.title
    title.text = "AV Preflight Test Deck"
    title.text_frame.paragraphs[0].runs[0].font.name = 'Georgia'
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].runs[0].font.bold = True

    subtitle = slide.placeholders[1]
    subtitle.text = "Generated by make_test_pptx.py"
    subtitle.text_frame.paragraphs[0].runs[0].font.name = 'Calibri'
    subtitle.text_frame.paragraphs[0].runs[0].font.size = Pt(20)

    slide.notes_slide.notes_text_frame.text = (
        "Welcome — this is the speaker notes for slide 1. "
        "Introduce the event and check that AV is set up correctly."
    )
    return slide


def slide2_image(prs):
    """Content slide — Arial + Gill Sans MT, embedded PNG, fade transition."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Slide with Embedded Image"
    title.text_frame.paragraphs[0].runs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)

    body = slide.placeholders[1]
    body.text = "This slide has a custom font and an embedded PNG image."
    body.text_frame.paragraphs[0].runs[0].font.name = 'Gill Sans MT'
    body.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    # Small PNG — intentionally tiny (80×60 px at 3"×2.25") to trigger blurry warning (~27 PPI)
    png_bytes = make_png(80, 60, rgb=(70, 130, 180))
    slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(5.5), Inches(2), Inches(3), Inches(2.25))

    # Large PNG — 1200×900 px displayed at 1.5"×1.125" = 800 PPI, triggers oversized warning
    large_png = make_png(1200, 900, rgb=(200, 80, 60))
    slide.shapes.add_picture(io.BytesIO(large_png), Inches(5.5), Inches(4.2), Inches(1.5), Inches(1.125))

    add_transition(slide, kind='fade', dur_ms=700)

    slide.notes_slide.notes_text_frame.text = "Slide 2 notes: check image is rendering correctly."
    return slide


def _add_video_pic(slide, shape_id, rId, x, y, w, h, fill_hex='1A1A2E'):
    """
    Inject a p:pic video shape into the slide's shape tree.
    Uses p14:media with r:link (linked video) in the nvPr extension.
    """
    P_NS   = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    A_NS   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    R_NS   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    P14_NS = 'http://schemas.microsoft.com/office/powerpoint/2010/main'

    pic_xml = f'''<p:pic
            xmlns:p="{P_NS}" xmlns:a="{A_NS}"
            xmlns:r="{R_NS}" xmlns:p14="{P14_NS}">
      <p:nvPicPr>
        <p:cNvPr id="{shape_id}" name="Video {shape_id}"/>
        <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
        <p:nvPr>
          <p:extLst>
            <p:ext uri="{{DAA4B4D4-6D71-4841-9C94-3317A9D6C41F}}">
              <p14:media r:link="{rId}"/>
            </p:ext>
          </p:extLst>
        </p:nvPr>
      </p:nvPicPr>
      <p:blipFill>
        <a:blip/><a:stretch><a:fillRect/></a:stretch>
      </p:blipFill>
      <p:spPr>
        <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="{fill_hex}"/></a:solidFill>
      </p:spPr>
    </p:pic>'''
    slide.shapes._spTree.append(etree.fromstring(pic_xml))


def _add_video_timing(slide, video_specs):
    """
    Build and inject a p:timing element for one or more video shapes.

    video_specs: list of dicts — {shape_id, autoplay, loop, muted}

    Trigger logic:
      autoplay=True  → <p:cond delay="0"/>   fires immediately on slide entry
      autoplay=False → <p:cond delay="indefinite"/>  waits for user click
    Loop:  repeatCount="indefinite" on the outer p:cTn
    Mute:  extra p:cmd setVolume(0) targeting the shape
    """
    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    sld = slide.element
    for old in sld.findall(qn('p:timing')):
        sld.remove(old)

    par_blocks = []
    for i, spec in enumerate(video_specs):
        sid     = str(spec['shape_id'])
        b       = 100 + i * 20
        trigger = '<p:cond delay="0"/>' if spec['autoplay'] else '<p:cond delay="indefinite"/>'
        repeat  = ' repeatCount="indefinite"' if spec['loop'] else ''
        mute_cmd = (
            f'<p:cmd type="call" cmd="setVolume(0)">'
            f'<p:cBhvr><p:cTn id="{b+3}" dur="1"/>'
            f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl></p:cBhvr></p:cmd>'
            if spec['muted'] else ''
        )
        par_blocks.append(f'''
        <p:par xmlns:p="{P_NS}">
          <p:cTn id="{b}" fill="hold"{repeat}>
            <p:stCondLst>{trigger}</p:stCondLst>
            <p:childTnLst>
              <p:par>
                <p:cTn id="{b+1}" fill="hold" nodeType="clickEffect">
                  <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                  <p:childTnLst>
                    <p:cmd type="call" cmd="playFrom(0.0)">
                      <p:cBhvr>
                        <p:cTn id="{b+2}" dur="10000"/>
                        <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
                      </p:cBhvr>
                    </p:cmd>
                    {mute_cmd}
                  </p:childTnLst>
                </p:cTn>
              </p:par>
            </p:childTnLst>
          </p:cTn>
        </p:par>''')

    timing_xml = f'''<p:timing xmlns:p="{P_NS}">
      <p:tnLst>
        <p:par>
          <p:cTn id="10" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="11" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>{''.join(par_blocks)}</p:childTnLst>
                </p:cTn>
                <p:prevCondLst><p:cond evt="onPrevNode" delay="0"/></p:prevCondLst>
                <p:nextCondLst><p:cond evt="onNextNode" delay="0"/></p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
    </p:timing>'''

    timing_elem = etree.fromstring(timing_xml)
    insert_at = len(sld)
    for i, child in enumerate(sld):
        if child.tag.split('}')[-1] == 'extLst':
            insert_at = i
            break
    sld.insert(insert_at, timing_elem)


def slide3_media(prs):
    """
    Media slide — two video shapes with contrasting playback settings + linked audio.

    Video 1 (left):  autoplay, not muted, no loop  → the ideal AV setup
    Video 2 (right): click-to-play, muted, looping → exercises all warning paths
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.75))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "Video Playback Settings Test"
    run.font.name = 'Trebuchet MS'
    run.font.size = Pt(28)

    # Labels
    for x_off, txt in ((0.5, "VIDEO 1 — autoplay · not muted · no loop"),
                        (5.2, "VIDEO 2 — click-to-play · muted · looping")):
        lb = slide.shapes.add_textbox(Inches(x_off), Inches(1.2), Inches(4.5), Inches(0.4))
        lr = lb.text_frame.paragraphs[0].add_run()
        lr.text = txt
        lr.font.name = 'Trebuchet MS'
        lr.font.size = Pt(11)

    # Assign shape IDs that won't clash with existing placeholder IDs
    max_id  = max(s.shape_id for s in slide.shapes)
    vid1_id = max_id + 1
    vid2_id = max_id + 2

    # Linked video relationships
    rId1 = slide.part.relate_to('C:/Videos/event_intro.mp4',  VIDEO_RELTYPE, is_external=True)
    rId2 = slide.part.relate_to('C:/Videos/bg_ambient.mp4',   VIDEO_RELTYPE, is_external=True)

    # Linked audio
    slide.part.relate_to('C:/Music/background.mp3', AUDIO_RELTYPE, is_external=True)

    # Video shapes (dark placeholder rectangles)
    _add_video_pic(slide, vid1_id, rId1,
                   int(Inches(0.5)), int(Inches(1.75)),
                   int(Inches(4.2)), int(Inches(2.4)), fill_hex='1A1A2E')
    _add_video_pic(slide, vid2_id, rId2,
                   int(Inches(5.2)), int(Inches(1.75)),
                   int(Inches(4.2)), int(Inches(2.4)), fill_hex='2E1A1A')

    # Timing: one autoplay, one click / loop / muted
    _add_video_timing(slide, [
        dict(shape_id=vid1_id, autoplay=True,  loop=False, muted=False),
        dict(shape_id=vid2_id, autoplay=False, loop=True,  muted=True),
    ])

    add_transition(slide, kind='push', dur_ms=500)
    return slide


def slide4_animations(prs):
    """Animation slide — Impact font, click-triggered fade-in, external hyperlink."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Title text box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8.5), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Animated Slide"
    run.font.name  = 'Impact'
    run.font.size  = Pt(40)
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Animated content box — fade in on click
    anim_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(2))
    tf2 = anim_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "This text box has a fade-in animation (click to reveal)."
    run2.font.name = 'Impact'
    run2.font.size = Pt(22)

    # Hyperlink box
    link_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(5), Inches(0.7))
    tf3 = link_box.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "Visit example.com (external hyperlink)"
    run3.font.name = 'Helvetica Neue'
    run3.font.size = Pt(16)
    run3.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
    run3.hyperlink.address = 'https://example.com'

    add_transition(slide, kind='dissolve', dur_ms=600)
    add_animation(slide, shape_id=anim_box.shape_id)
    return slide


def slide5_hidden(prs):
    """Hidden slide — Courier New, no notes. Tests hidden slide detection."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Hidden Slide (should not show)"
    title.text_frame.paragraphs[0].runs[0].font.name = 'Courier New'
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

    body = slide.placeholders[1]
    body.text = "This slide is hidden. The preflight tool should flag it."
    body.text_frame.paragraphs[0].runs[0].font.name = 'Courier New'

    # Mark as hidden
    slide.element.set('show', '0')
    return slide


# ── Build ─────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Emu(9144000)   # 10 inches — 16:9
    prs.slide_height = Emu(5143500)   # 5.625 inches

    slide1_title(prs)
    slide2_image(prs)
    slide3_media(prs)
    slide4_animations(prs)
    slide5_hidden(prs)

    prs.save(OUT)
    print(f"Written: {OUT}")
    print(f"  5 slides")
    print(f"  Fonts:       Georgia, Calibri, Arial, Gill Sans MT, Trebuchet MS, Impact, Helvetica Neue, Courier New")
    print(f"  Media:       2 linked videos (autoplay+normal; loop+mute), 1 linked audio")
    print(f"  Images:      2 embedded PNGs (1 blurry ~27 PPI, 1 oversized ~800 PPI)")
    print(f"  Transitions: fade (slide 2), push+auto (slide 3), dissolve (slide 4)")
    print(f"  Animations:  1 animEffect on slide 4")
    print(f"  Hyperlinks:  1 external (slide 4)")
    print(f"  Notes:       slides 1 + 2")
    print(f"  Hidden:      slide 5")


if __name__ == '__main__':
    build()
