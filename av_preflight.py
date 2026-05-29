#!/usr/bin/env python3
"""
AV Preflight Tool for PowerPoint presentations.

Analyzes a .pptx file and reports on everything that matters before
taking it to a live event: fonts, embedded vs linked media, transitions,
animations, hidden slides, hyperlinks, and embedded objects.

Usage:
    python av_preflight.py presentation.pptx
    python av_preflight.py presentation.pptx --display 3840x2160

The --display flag defaults to 1920x1080. The deck's aspect ratio is
compared against the target display and a mismatch is flagged.

Requires:
    pip install python-pptx
"""

import argparse
import os
import shutil
import struct
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path, PurePosixPath
from collections import defaultdict, Counter

try:
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# ── Constants ─────────────────────────────────────────────────────────────────

EMU_PER_INCH = 914400

VIDEO_EXTENSIONS  = {'.mp4', '.mov', '.avi', '.wmv', '.mpg', '.mpeg',
                     '.m4v', '.mkv', '.webm', '.flv'}
AUDIO_EXTENSIONS  = {'.mp3', '.wav', '.aac', '.m4a', '.ogg', '.wma', '.flac'}
IMAGE_EXTENSIONS  = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff',
                     '.tif', '.emf', '.wmf', '.svg', '.ico'}

SAFE_VIDEO_FORMATS = {'.mp4', '.mov'}  # broadly supported on event systems
SAFE_AUDIO_FORMATS = {'.mp3', '.wav', '.aac', '.m4a'}

TRANS_WARN_MS = 1500  # transition duration above this is flagged as slow
ANIM_WARN_MS  = 1500  # animation effect duration above this is flagged as slow


# ── Formatting helpers ────────────────────────────────────────────────────────

def fmt_ms(ms):
    """Format milliseconds as a seconds string: 700 → '0.7s', 2500 → '2.5s'."""
    return f"{int(ms) / 1000:.1f}s"


def fmt_size(n):
    if n is None:
        return "?"
    for unit in ('B', 'KB', 'MB', 'GB'):
        if n < 1024:
            return f"{n:.0f} {unit}" if unit == 'B' else f"{n:.1f} {unit}"
        n /= 1024
    return f"{n:.1f} TB"


def aspect_label(w_emu, h_emu):
    r = w_emu / h_emu
    for val, name in ((16/9, "16:9"), (4/3, "4:3"), (16/10, "16:10"), (1.0, "1:1")):
        if abs(r - val) < 0.02:
            return name
    return f"custom ({r:.3f}:1)"


def parse_display(value):
    """Parse 'WxH' or 'W×H' into (width, height) integers."""
    for sep in ('x', '×', 'X'):
        if sep in value:
            w, h = value.split(sep, 1)
            return int(w.strip()), int(h.strip())
    raise argparse.ArgumentTypeError(f"expected WxH format (e.g. 1920x1080), got: {value!r}")


def aspect_match(deck_w_emu, deck_h_emu, disp_w, disp_h, tol=0.02):
    """Return (match, deck_ratio, display_ratio)."""
    deck_r = deck_w_emu / deck_h_emu
    disp_r = disp_w / disp_h
    return abs(deck_r - disp_r) < tol, deck_r, disp_r


def fmt_slide_list(slide_nums, show_count=True):
    """
    Format a list of slide numbers with run-length range compression.

      [1, 2, 3, 5, 7, 8, 9]  →  "[1–3, 5, 7–9]"
      [8, 11, 13, 14, 48, …, 103]  →  "[8, 11, 13–14, 48–103]  (60 slides)"

    A count in parentheses is appended when > 8 items, unless show_count=False
    (use False when the count is already stated in the surrounding text).
    """
    if not slide_nums:
        return "—"
    nums = sorted(int(n) for n in slide_nums)
    n    = len(nums)

    # Build consecutive runs
    runs = []
    s = e = nums[0]
    for v in nums[1:]:
        if v == e + 1:
            e = v
        else:
            runs.append((s, e))
            s = e = v
    runs.append((s, e))

    parts = [f"{s}–{e}" if e > s else str(s) for s, e in runs]
    result = "[" + ", ".join(parts) + "]"
    if n > 8 and show_count:
        result += f"  ({n} slides)"
    return result


def section_header(title):
    print(f"\n{title}")
    print("─" * 60)


# ── Font collection ───────────────────────────────────────────────────────────

def collect_fonts(element, slide_num, bucket):
    """
    Walk element tree collecting explicit font names from a:latin/@typeface.
    Theme font placeholders (+mj-lt, +mn-lt) are skipped here and reported
    separately via theme_font_names().
    """
    for latin in element.iter(qn('a:latin')):
        tf = latin.get('typeface', '').strip()
        if tf and not tf.startswith('+'):
            bucket[tf].add(slide_num)


def theme_font_names(prs):
    """Return (major_font, minor_font) from the slide master's theme."""
    names = {}
    try:
        master = prs.slide_master
        theme_part = master.part.part_related_by(
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme'
        )
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        fs = theme_part.element.find(f'{{{ns}}}themeElements/{{{ns}}}fontScheme')
        if fs is not None:
            for role in ('majorFont', 'minorFont'):
                node = fs.find(f'{{{ns}}}{role}')
                if node is not None:
                    latin = node.find(f'{{{ns}}}latin')
                    if latin is not None:
                        names[role] = latin.get('typeface', '')
    except Exception:
        pass
    return names


# ── Document properties ───────────────────────────────────────────────────────

# AppVersion major number → Office release label
_OFFICE_VERSION = {
    '9':  'Office 2000',
    '10': 'Office XP (2002)',
    '11': 'Office 2003',
    '12': 'Office 2007',
    '14': 'Office 2010',
    '15': 'Office 2013',
    '16': 'Office 2016 / 2019 / 2021 / 365',
}


def _office_version_label(app_version_str):
    """'16.0000' → 'Office 2016 / 2019 / 2021 / 365'"""
    if not app_version_str:
        return None
    major = app_version_str.split('.')[0]
    return _OFFICE_VERSION.get(major, f'v{app_version_str}')


def _os_from_app(app_str):
    """
    Infer authoring OS from the Application string.
    Returns 'Windows', 'macOS', 'Linux', or None.
    """
    if not app_str:
        return None
    low = app_str.lower()
    # PowerPoint for Mac explicitly says 'macintosh'
    if 'macintosh' in low:
        return 'macOS'
    # LibreOffice/Impress embeds the full platform: "LibreOffice/7.x$Linux_X86_64…"
    if 'linux' in low:
        return 'Linux'
    if 'macosx' in low or 'darwin' in low:
        return 'macOS'
    # Plain "Microsoft Office PowerPoint" without 'macintosh' → Windows
    if 'microsoft' in low:
        return 'Windows'
    return None


def document_properties(zf):
    """
    Parse docProps/core.xml and docProps/app.xml from the zip.
    Returns a dict with keys:
      title, author, last_modified_by, created, modified, revision,
      app_name, app_version, office_label, os_hint, company, presentation_format
    All values may be None if absent.
    """
    import xml.etree.ElementTree as ET

    CP  = 'http://schemas.openxmlformats.org/package/2006/metadata/core-properties'
    DC  = 'http://purl.org/dc/elements/1.1/'
    DCT = 'http://purl.org/dc/terms/'
    APP = 'http://schemas.openxmlformats.org/officeDocument/2006/extended-properties'

    props = dict(
        title=None, author=None, last_modified_by=None,
        created=None, modified=None, revision=None,
        app_name=None, app_version=None, office_label=None,
        os_hint=None, company=None, presentation_format=None,
    )

    # ── core.xml ──────────────────────────────────────────────────────────────
    try:
        raw = zf.read('docProps/core.xml')
        root = ET.fromstring(raw)

        def _text(tag_ns, tag_local):
            el = root.find(f'{{{tag_ns}}}{tag_local}')
            return el.text.strip() if el is not None and el.text else None

        props['title']            = _text(DC,  'title')
        props['author']           = _text(DC,  'creator')
        props['last_modified_by'] = _text(CP,  'lastModifiedBy')
        props['created']          = _text(DCT, 'created')
        props['modified']         = _text(DCT, 'modified')
        props['revision']         = _text(CP,  'revision')
    except Exception:
        pass

    # ── app.xml ───────────────────────────────────────────────────────────────
    try:
        raw = zf.read('docProps/app.xml')
        root = ET.fromstring(raw)

        def _atext(tag_local):
            el = root.find(f'{{{APP}}}{tag_local}')
            return el.text.strip() if el is not None and el.text else None

        props['app_name']            = _atext('Application')
        props['app_version']         = _atext('AppVersion')
        props['company']             = _atext('Company')
        props['presentation_format'] = _atext('PresentationFormat')
    except Exception:
        pass

    props['office_label'] = _office_version_label(props['app_version'])
    props['os_hint']      = _os_from_app(props['app_name'])

    return props


def _fmt_datetime(iso_str):
    """'2024-05-28T14:22:00Z' → '2024-05-28  14:22 UTC'"""
    if not iso_str:
        return None
    try:
        date, time = iso_str.rstrip('Z').split('T')
        return f"{date}  {time[:5]} UTC"
    except Exception:
        return iso_str


# ── Media extraction ──────────────────────────────────────────────────────────

def _embedded_size(rel, zf):
    try:
        zip_path = str(rel.target_part.partname).lstrip('/')
        return zf.getinfo(zip_path).file_size
    except Exception:
        return None


def slide_media(slide, slide_num, zf):
    """Return list of video/audio items from a slide's relationships."""
    items = []
    for rel in slide.part.rels.values():
        rt = rel.reltype.lower()
        target = rel.target_ref
        ext = PurePosixPath(target).suffix.lower()

        if 'video' in rt or ext in VIDEO_EXTENSIONS:
            kind = 'video'
        elif 'audio' in rt or ext in AUDIO_EXTENSIONS:
            kind = 'audio'
        else:
            continue

        embedded = not rel.is_external
        items.append(dict(
            slide=slide_num,
            kind=kind,
            name=PurePosixPath(target).name,
            ext=ext,
            embedded=embedded,
            linked_path=target if rel.is_external else None,
            size=_embedded_size(rel, zf) if embedded else None,
        ))
    return items


def slide_images(slide, slide_num, zf):
    items = []
    for rel in slide.part.rels.values():
        if 'image' not in rel.reltype.lower():
            continue
        target = rel.target_ref
        ext = PurePosixPath(target).suffix.lower()
        embedded = not rel.is_external
        items.append(dict(
            slide=slide_num,
            name=PurePosixPath(target).name,
            ext=ext,
            embedded=embedded,
            linked_path=target if rel.is_external else None,
            size=_embedded_size(rel, zf) if embedded else None,
        ))
    return items


# ── Image resolution check ───────────────────────────────────────────────────

# PPI thresholds for screen display
_PPI_LOW    = 96   # below this may look blurry on a standard monitor
_PPI_HIGH   = 300  # above this is print resolution — no visible benefit on screen
_PPI_TARGET = 192  # used to suggest ideal pixel dimensions (good for HiDPI/Retina)

R_NS   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
P14_NS = 'http://schemas.microsoft.com/office/powerpoint/2010/main'


def image_ppi_issues(prs):
    """
    Inspect every picture shape for under- or over-resolution.
    Returns (issues, n_checked) where issues is None if Pillow is not installed.
    """
    try:
        import io
        from PIL import Image as PILImage
    except ImportError:
        return None, 0

    issues  = []
    checked = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob   = shape.image.blob
                img    = PILImage.open(io.BytesIO(blob))
                img_w, img_h = img.size
                checked += 1

                if shape.width == 0 or shape.height == 0:
                    continue

                shape_w_in = shape.width  / EMU_PER_INCH
                shape_h_in = shape.height / EMU_PER_INCH
                ppi = ((img_w / shape_w_in) + (img_h / shape_h_in)) / 2

                if _PPI_LOW <= ppi <= _PPI_HIGH:
                    continue  # fine — skip

                # Resolve filename from the blip relationship
                try:
                    blip = shape._element.find('.//' + qn('a:blip'))
                    rId  = blip.get(f'{{{R_NS}}}embed')
                    name = Path(str(shape.part.rels[rId].target_part.partname)).name
                except Exception:
                    name = f'<shape {shape.shape_id}>'

                ideal_w = round(shape_w_in * _PPI_TARGET)
                ideal_h = round(shape_h_in * _PPI_TARGET)

                issues.append(dict(
                    slide      = slide_num,
                    name       = name,
                    img_w      = img_w,
                    img_h      = img_h,
                    w_in       = round(shape_w_in, 2),
                    h_in       = round(shape_h_in, 2),
                    ppi        = round(ppi),
                    file_size  = len(blob),
                    too_low    = ppi < _PPI_LOW,
                    too_high   = ppi > _PPI_HIGH,
                    ideal_w    = ideal_w,
                    ideal_h    = ideal_h,
                    # How many times bigger (linear) than the HiDPI target
                    oversize_x = round(ppi / _PPI_TARGET, 1) if ppi > _PPI_HIGH else None,
                ))
            except Exception:
                continue

    return issues, checked


# ── Video shape detection ─────────────────────────────────────────────────────

def get_video_playback_props(slide, shape_id):
    """
    Parse the slide's p:timing to determine autoplay / loop / muted for a video shape.

    PowerPoint stores media triggers inside the mainSeq timeline.  The OUTER
    p:par (direct child of mainSeq's childTnLst) controls when playback starts:
      • stCondLst/cond delay="0" with no evt → fires on slide entry → autoplay
      • stCondLst/cond delay="indefinite"    → waits for click  → not autoplay

    repeatCount="indefinite" on that outer cTn = loop.
    A p:cmd type="call" cmd="setVolume(0)" targeting the shape = muted.

    Returns dict(autoplay, loop, muted), or None if no timing block found.
    """
    sid    = str(shape_id)
    timing = slide.element.find(qn('p:timing'))
    if timing is None:
        return None

    # Find the mainSeq container node
    main_cTn = None
    for cTn in timing.findall('.//' + qn('p:cTn')):
        if cTn.get('nodeType') == 'mainSeq':
            main_cTn = cTn
            break
    if main_cTn is None:
        return None

    child_lst = main_cTn.find(qn('p:childTnLst'))
    if child_lst is None:
        return None

    for outer_par in child_lst.findall(qn('p:par')):
        # Check whether this block references our video shape at all
        if not any(sp.get('spid') == sid
                   for sp in outer_par.findall('.//' + qn('p:spTgt'))):
            continue

        outer_cTn = outer_par.find(qn('p:cTn'))
        if outer_cTn is None:
            continue

        # Autoplay: outer stCondLst has <p:cond delay="0"> with NO evt attribute
        autoplay = False
        stCond = outer_cTn.find(qn('p:stCondLst'))
        if stCond is not None:
            for cond in stCond.findall(qn('p:cond')):
                if cond.get('delay') == '0' and cond.get('evt') is None:
                    autoplay = True
                    break

        # Loop: repeatCount="indefinite" on the outer cTn
        loop = outer_cTn.get('repeatCount') == 'indefinite'

        # Mute: any p:cmd setVolume(0) targeting our shape
        muted = False
        for cmd in outer_par.findall('.//' + qn('p:cmd')):
            if cmd.get('type') == 'call' and cmd.get('cmd') == 'setVolume(0)':
                if any(sp.get('spid') == sid
                       for sp in cmd.findall('.//' + qn('p:spTgt'))):
                    muted = True
                    break

        return dict(autoplay=autoplay, loop=loop, muted=muted)

    return None  # shape referenced in no timing block


# ── Video codec detection ─────────────────────────────────────────────────────

_CODEC_LABELS = {
    'avc1': 'H.264',   'avc3': 'H.264',
    'hvc1': 'H.265',   'hev1': 'H.265',
    'av01': 'AV1',
    'vp08': 'VP8',     'vp09': 'VP9',
    'apch': 'ProRes',  'apcn': 'ProRes',  'apcs': 'ProRes',
    'apco': 'ProRes',  'ap4h': 'ProRes',  'ap4x': 'ProRes',
    'dvh1': 'Dolby Vision',  'dvhe': 'Dolby Vision',
    'mp4v': 'MPEG-4',
    'mjp2': 'Motion JPEG',   'jpeg': 'Motion JPEG',
    'h263': 'H.263',   's263': 'H.263',
}

_CODEC_SAFE = {'H.264', 'MPEG-4'}

_CODEC_WARN = {
    'H.265':         'hardware decode required — may fail on older event systems',
    'ProRes':        'Apple codec — will NOT play on Windows without QuickTime or a codec pack',
    'AV1':           'limited support — not recommended for live events',
    'VP8':           'web codec — unlikely to play in PowerPoint on Windows',
    'VP9':           'web codec — unlikely to play in PowerPoint on Windows',
    'Dolby Vision':  'HDR — may not render correctly on standard displays',
    'Motion JPEG':   'large file size — verify playback on event system',
    'H.263':         'legacy codec — verify playback on event system',
    'WMV':           'Windows-only — will not play on macOS without a codec',
    'AVI':           'container only — codec unknown; verify playback on event system',
    'WebM/MKV':      'open web format — not reliably supported in PowerPoint',
}

_CODEC_PROBE_BYTES = 256 * 1024  # read first 256 KB — enough for moov at start of file


def _probe_mp4_codec(data: bytes):
    """Parse ISOBMFF (MP4/MOV) bytes and return the video codec label or raw fourcc."""
    def _iter(buf, start, end):
        pos = start
        while pos + 8 <= end:
            sz = struct.unpack_from('>I', buf, pos)[0]
            try:
                bt = buf[pos+4:pos+8].decode('ascii')
            except Exception:
                return
            if sz == 1:                   # 64-bit extended size
                if pos + 16 > end:
                    return
                sz  = struct.unpack_from('>Q', buf, pos+8)[0]
                hdr = 16
            elif sz == 0:                 # box runs to EOF
                sz  = end - pos
                hdr = 8
            else:
                hdr = 8
            if sz < hdr or pos + sz > end:
                return
            yield bt, pos + hdr, pos + sz
            pos += sz

    def _find(buf, start, end, *path):
        cs, ce = start, end
        for name in path:
            for bt, bs, be in _iter(buf, cs, ce):
                if bt == name:
                    cs, ce = bs, be
                    break
            else:
                return None, None
        return cs, ce

    n = len(data)
    # Locate moov box
    moov_s = moov_e = None
    for bt, bs, be in _iter(data, 0, n):
        if bt == 'moov':
            moov_s, moov_e = bs, be
            break
    if moov_s is None:
        return None   # moov not in first 256 KB (non-fast-start file)

    # Walk trak boxes looking for the video track
    for bt, bs, be in _iter(data, moov_s, moov_e):
        if bt != 'trak':
            continue
        mdia_s, mdia_e = _find(data, bs, be, 'mdia')
        if mdia_s is None:
            continue
        # hdlr handler_type must be 'vide'
        hdlr_s, _ = _find(data, mdia_s, mdia_e, 'hdlr')
        if hdlr_s is None or hdlr_s + 12 > n:
            continue
        # hdlr content: version(1) + flags(3) + pre_defined(4) + handler_type(4)
        if data[hdlr_s+8:hdlr_s+12] != b'vide':
            continue
        # Navigate to stsd (sample description — contains codec fourcc)
        stsd_s, stsd_e = _find(data, mdia_s, mdia_e, 'minf', 'stbl', 'stsd')
        if stsd_s is None:
            continue
        # stsd content: version(1) + flags(3) + entry_count(4) + first entry...
        entry = stsd_s + 8
        if entry + 8 > stsd_e:
            continue
        # First sample entry: size(4) + codec_fourcc(4)
        try:
            fourcc = data[entry+4:entry+8].decode('ascii').strip()
            return _CODEC_LABELS.get(fourcc, fourcc)
        except Exception:
            continue
    return None


# ffprobe codec names → our labels
_FFPROBE_CODEC_MAP = {
    'h264':        'H.264',   'h265':  'H.265',  'hevc':   'H.265',
    'prores':      'ProRes',  'av1':   'AV1',
    'vp8':         'VP8',     'vp9':   'VP9',
    'mpeg4':       'MPEG-4',  'mpeg2video': 'MPEG-2',
    'wmv1':        'WMV',     'wmv2':  'WMV',    'wmv3':   'WMV',   'vc1': 'WMV',
    'dvvideo':     'DV',      'mjpeg': 'Motion JPEG',
    'theora':      'Theora',
}

_FFPROBE_SIZE_LIMIT = 500 * 1024 * 1024  # skip extraction for files > 500 MB
_ffprobe_exe = None   # lazily resolved: False = not found, str = path


def _get_ffprobe():
    global _ffprobe_exe
    if _ffprobe_exe is None:
        _ffprobe_exe = shutil.which('ffprobe') or False
    return _ffprobe_exe or None


def _ffprobe_codec(rel, zf, ext: str):
    """Extract video to a temp file and run ffprobe. Returns codec label or None."""
    fp = _get_ffprobe()
    if not fp:
        return None
    try:
        zip_path = str(rel.target_part.partname).lstrip('/')
        file_size = zf.getinfo(zip_path).file_size
        if file_size > _FFPROBE_SIZE_LIMIT:
            return None   # too large to extract
        suffix = ext if ext.startswith('.') else f'.{ext}'
        fd, tmp_path = tempfile.mkstemp(suffix=suffix)
        try:
            with os.fdopen(fd, 'wb') as fout:
                fout.write(zf.read(zip_path))
            result = subprocess.run(
                [fp, '-v', 'error',
                 '-select_streams', 'v:0',
                 '-show_entries', 'stream=codec_name',
                 '-of', 'default=noprint_wrappers=1:nokey=1',
                 tmp_path],
                capture_output=True, text=True, timeout=15,
            )
            codec_raw = result.stdout.strip().lower()
            if codec_raw:
                return _FFPROBE_CODEC_MAP.get(codec_raw, codec_raw.upper())
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass
    except Exception:
        pass
    return None


def _video_codec(rel, zf, ext: str):
    """Detect codec for an embedded video. Returns a label string or None."""
    try:
        if ext in ('.mp4', '.mov', '.m4v'):
            # Fast path: custom ISOBMFF parser on first 256 KB (handles fast-start files)
            zip_path = str(rel.target_part.partname).lstrip('/')
            with zf.open(zip_path) as f:
                data = f.read(_CODEC_PROBE_BYTES)
            result = _probe_mp4_codec(data)
            if result is not None:
                return result
            # moov not in first 256 KB (non-fast-start) — fall through to ffprobe
        # ffprobe handles all formats and moov-at-end files
        fp_result = _ffprobe_codec(rel, zf, ext)
        if fp_result:
            return fp_result
        # Container-type fallback when ffprobe is unavailable
        if ext == '.wmv':
            return 'WMV'
        if ext == '.avi':
            return 'AVI'
        if ext in ('.mkv', '.webm'):
            return 'WebM/MKV'
    except Exception:
        pass
    return None


def find_video_shapes(slide, slide_num, zf):
    """
    Detect video shapes by looking for the p14:media extension inside p:nvPr.
    Returns a list of video dicts, each containing '_target' for deduplication.
    Playback properties are resolved via get_video_playback_props().

    Note: modern PPTX stores each embedded video in TWO relationships — a
    'media' reltype (used by p14:media) and a legacy 'video' reltype.  Both
    point to the same file, so deduplication must be by target URL, not rId.
    """
    results = []
    for shape in slide.shapes:
        nvPr = shape.element.find('.//' + qn('p:nvPr'))
        if nvPr is None:
            continue

        media_el = nvPr.find(f'.//{{{P14_NS}}}media')
        if media_el is None:
            continue

        rId_link  = media_el.get(f'{{{R_NS}}}link')
        rId_embed = media_el.get(f'{{{R_NS}}}embed')
        rId       = rId_link or rId_embed
        embedded  = rId_embed is not None

        if rId is None:
            continue

        try:
            rel    = slide.part.rels[rId]
            target = rel.target_ref
            ext    = PurePosixPath(target).suffix.lower()
            if ext not in VIDEO_EXTENSIONS:
                continue
            name  = PurePosixPath(target).name
            size  = _embedded_size(rel, zf) if embedded else None
            codec = _video_codec(rel, zf, ext) if embedded else None
        except Exception:
            continue

        props = get_video_playback_props(slide, shape.shape_id)

        results.append(dict(
            slide       = slide_num,
            kind        = 'video',
            name        = name,
            ext         = ext,
            embedded    = embedded,
            linked_path = target if not embedded else None,
            size        = size,
            codec       = codec,
            autoplay    = props['autoplay'] if props else None,
            loop        = props['loop']     if props else None,
            muted       = props['muted']    if props else None,
            _target     = target,            # internal — stripped in analyze()
        ))

    return results


# ── Transitions & animations ──────────────────────────────────────────────────

def slide_transition(slide, n):
    t = slide.element.find('.//' + qn('p:transition'))
    if t is None:
        return None
    children = [c.tag.split('}')[-1] for c in t if '}' in c.tag]
    raw = t.get('dur')
    try:
        dur_ms = int(raw) if raw else None
    except ValueError:
        dur_ms = None
    return dict(
        slide=n,
        kind=children[0] if children else 'none',
        click=t.get('advClick', '1') != '0',
        auto_ms=t.get('advTm'),
        dur_ms=dur_ms,        # int ms, or None → PowerPoint default (~700 ms)
    )


_ANIM_TAGS = [qn(t) for t in
              ('p:anim', 'p:animEffect', 'p:animMotion', 'p:animScale', 'p:animRot')]


def animation_effect_durations(element):
    """
    Return a list of per-effect durations (int ms, or None for 'indefinite')
    for any element that can hold p:timing — slide, master, or layout.
    """
    timing = element.find(qn('p:timing'))
    if timing is None:
        return []

    durations = []
    for tag in _ANIM_TAGS:
        for anim in timing.findall('.//' + tag):
            cBhvr = anim.find(qn('p:cBhvr'))
            if cBhvr is None:
                continue
            cTn = cBhvr.find(qn('p:cTn'))
            if cTn is None:
                continue
            raw = cTn.get('dur', '')
            if raw == 'indefinite':
                durations.append(None)
            elif raw:
                try:
                    durations.append(int(raw))
                except ValueError:
                    pass
    return durations


def check_masters_for_animations(prs):
    """
    Check the slide master and every slide layout for animation effects.
    Returns list of (label, durations) for any that have animations.
    """
    found = []
    master = prs.slide_master
    durs = animation_effect_durations(master.element)
    if durs:
        found.append(('Slide master', durs))
    for layout in master.slide_layouts:
        durs = animation_effect_durations(layout.element)
        if durs:
            found.append((f"Layout '{layout.name or 'unnamed'}'", durs))
    return found


# ── Presentation-level properties ─────────────────────────────────────────────

def show_properties(prs):
    """Extract playback mode and loop setting from p:showPr."""
    props = {'mode': 'standard', 'loop': False, 'no_narration': False}
    showPr = prs.element.find(qn('p:showPr'))
    if showPr is None:
        return props
    props['loop'] = showPr.get('loop', '0') == '1'
    props['no_narration'] = showPr.get('showNarration', '1') == '0'
    if showPr.find(qn('p:kiosk')) is not None:
        props['mode'] = 'kiosk (auto-advance, no menu)'
    elif showPr.find(qn('p:browse')) is not None:
        props['mode'] = 'browsed in window'
    else:
        props['mode'] = 'presented by speaker'
    return props


def presentation_sections(prs):
    """Return list of section names if the deck uses sections."""
    sections = []
    sectionLst = prs.element.find('.//' + qn('p:sectionLst'))
    if sectionLst is not None:
        for sec in sectionLst.findall(qn('p:section')):
            name = sec.get('name', '(unnamed)')
            sections.append(name)
    return sections


# ── Main ──────────────────────────────────────────────────────────────────────

def analyze(path, display_wh=(1920, 1080), verbose=False):
    path = Path(path)
    if not path.exists():
        print(f"File not found: {path}")
        sys.exit(1)

    prs = Presentation(str(path))
    zf  = zipfile.ZipFile(str(path))

    print(f"\n{'━'*60}")
    print(f"  AV PREFLIGHT — {path.name}")
    print(f"{'━'*60}")

    # ── Deck overview ──────────────────────────────────────────────────────────
    n_slides = len(prs.slides)
    w_emu, h_emu = prs.slide_width, prs.slide_height
    w_in, h_in   = w_emu / EMU_PER_INCH, h_emu / EMU_PER_INCH
    w_px, h_px   = round(w_in * 96), round(h_in * 96)

    hidden_slides = [i+1 for i, s in enumerate(prs.slides)
                     if s.element.get('show') == '0']
    show_props = show_properties(prs)
    sections   = presentation_sections(prs)

    section_header("DECK OVERVIEW")
    print(f"  Slides          {n_slides}" +
          (f"  ({len(hidden_slides)} hidden: {fmt_slide_list(hidden_slides, show_count=False)})" if hidden_slides else ""))
    print(f"  Dimensions      {w_in:.2f}\" × {h_in:.2f}\"  "
          f"({w_px} × {h_px} px @ 96 dpi)")
    print(f"  Aspect ratio    {aspect_label(w_emu, h_emu)}")

    disp_w, disp_h = display_wh
    ratio_ok, deck_r, disp_r = aspect_match(w_emu, h_emu, disp_w, disp_h)
    disp_label = f"{disp_w}×{disp_h}  ({aspect_label(disp_w, disp_h)})"
    if ratio_ok:
        match_str = "✓ matches"
    else:
        # Describe the visual consequence
        if deck_r < disp_r:
            consequence = "pillarboxed (black bars left/right)"
        else:
            consequence = "letterboxed (black bars top/bottom)"
        match_str = f"✗ MISMATCH — {consequence}"
    print(f"  Target display  {disp_label}")
    print(f"  Ratio match     {match_str}")

    print(f"  File size       {fmt_size(path.stat().st_size)}")
    print(f"  Show mode       {show_props['mode']}")
    if show_props['loop']:
        print(f"  Loop            YES")
    if sections:
        print(f"  Sections        {len(sections)}: {', '.join(sections)}")

    # ── Document properties ────────────────────────────────────────────────────
    doc = document_properties(zf)

    section_header("DOCUMENT PROPERTIES")
    if doc['title']:
        print(f"  Title           {doc['title']}")
    if doc['author']:
        print(f"  Author          {doc['author']}")
    if doc['last_modified_by'] and doc['last_modified_by'] != doc['author']:
        print(f"  Last saved by   {doc['last_modified_by']}")
    if doc['created']:
        print(f"  Created         {_fmt_datetime(doc['created'])}")
    if doc['modified']:
        print(f"  Modified        {_fmt_datetime(doc['modified'])}")
    if doc['revision']:
        print(f"  Revision        {doc['revision']}")

    # Application / authoring environment
    app_str = doc['app_name'] or '?'
    if doc['office_label']:
        app_str += f"  ({doc['office_label']})"
    if doc['os_hint']:
        app_str += f"  [{doc['os_hint']}]"
    print(f"  Application     {app_str}")

    if doc['company']:
        print(f"  Company         {doc['company']}")
    if doc['presentation_format']:
        print(f"  Format          {doc['presentation_format']}")

    # Warn if the authoring app is not PowerPoint (LibreOffice, Keynote, etc.)
    if doc['app_name'] and 'microsoft' not in doc['app_name'].lower():
        print(f"\n  ⚠  Not authored in Microsoft PowerPoint — "
              f"verify rendering on the event system")

    # ── Fonts ──────────────────────────────────────────────────────────────────
    fonts = defaultdict(set)
    collect_fonts(prs.slide_master.element, 0, fonts)
    for i, slide in enumerate(prs.slides, 1):
        collect_fonts(slide.element, i, fonts)

    theme_fonts  = theme_font_names(prs)
    # Embedded font files stored inside the .pptx zip
    embedded_font_files = [f for f in zf.namelist() if f.startswith('ppt/fonts/')]

    section_header(f"FONTS  ({len(fonts)} explicit)")
    for name, slides in sorted(fonts.items()):
        on_slides  = sorted(s for s in slides if s > 0)
        tags       = []
        if 0 in slides:
            tags.append("master")
        if name == theme_fonts.get('majorFont'):
            tags.append("theme heading")
        if name == theme_fonts.get('minorFont'):
            tags.append("theme body")
        tag_str = f"  [{', '.join(tags)}]" if tags else ""
        print(f"  {name:<40} slides {fmt_slide_list(on_slides)}{tag_str}")

    if theme_fonts:
        print(f"\n  Theme heading font  {theme_fonts.get('majorFont', '?')}")
        print(f"  Theme body font     {theme_fonts.get('minorFont', '?')}")

    if embedded_font_files:
        print(f"\n  Embedded font files ({len(embedded_font_files)}):")
        for f in embedded_font_files:
            print(f"    {Path(f).name}  ({fmt_size(zf.getinfo(f).file_size)})")
    else:
        print(f"\n  ⚠  No embedded fonts — all fonts must be installed on "
              f"the playback machine")

    # ── Media ──────────────────────────────────────────────────────────────────
    all_videos = []
    all_audios = []
    all_images = []

    for i, slide in enumerate(prs.slides, 1):
        # Shape-based detection gives autoplay / loop / muted.
        # Dedup by target URL: modern PPTX stores each embedded video in both
        # a 'media' reltype (used by p14:media) AND a legacy 'video' reltype —
        # same file, different rIds, so we match on the resolved target path.
        shape_vids    = find_video_shapes(slide, i, zf)
        shape_targets = {v.pop('_target') for v in shape_vids}
        all_videos.extend(shape_vids)

        # Relationship scan: audio + any videos not already found via shapes
        for rel in slide.part.rels.values():
            rt     = rel.reltype.lower()
            target = rel.target_ref
            ext    = PurePosixPath(target).suffix.lower()

            if 'audio' in rt or ext in AUDIO_EXTENSIONS:
                embedded = not rel.is_external
                all_audios.append(dict(
                    slide       = i,
                    kind        = 'audio',
                    name        = PurePosixPath(target).name,
                    ext         = ext,
                    embedded    = embedded,
                    linked_path = target if rel.is_external else None,
                    size        = _embedded_size(rel, zf) if embedded else None,
                ))
            elif ('video' in rt or ext in VIDEO_EXTENSIONS) and target not in shape_targets:
                # Video rel with no matching shape (legacy deck or bare relationship)
                embedded = not rel.is_external
                all_videos.append(dict(
                    slide       = i,
                    kind        = 'video',
                    name        = PurePosixPath(target).name,
                    ext         = ext,
                    embedded    = embedded,
                    linked_path = target if rel.is_external else None,
                    size        = _embedded_size(rel, zf) if embedded else None,
                    codec       = _video_codec(rel, zf, ext) if embedded else None,
                    autoplay    = None,
                    loop        = None,
                    muted       = None,
                ))

        all_images.extend(slide_images(slide, i, zf))

    # Post-collection dedup: same slide+name can arrive from both the shape scan
    # AND the legacy 'video' rel (if target_ref strings differ between rel types).
    # Keep the entry that has playback data (autoplay is not None) over the one
    # that is all-None; otherwise keep first seen.
    _seen_vids = {}
    _deduped   = []
    for v in all_videos:
        key = (v['slide'], v['name'])
        if key not in _seen_vids:
            _seen_vids[key] = len(_deduped)
            _deduped.append(v)
        elif v['autoplay'] is not None and _deduped[_seen_vids[key]]['autoplay'] is None:
            _deduped[_seen_vids[key]] = v   # upgrade to the one with playback props
    all_videos = _deduped

    # ── VIDEO section — chart-style table ──────────────────────────────────────
    section_header(f"VIDEO  ({len(all_videos)})")
    if not all_videos:
        print("  (none)")
    else:
        print(f"  {'Slide':>5}  {'File':<30}  {'Source':<9}  {'Codec':<14}"
              f"  {'Autoplay':<8}  {'Loop':<5}  {'Mute':<5}  Size")
        print(f"  {'─'*5}  {'─'*30}  {'─'*9}  {'─'*14}"
              f"  {'─'*8}  {'─'*5}  {'─'*5}  {'─'*10}")
        has_unknown = any(v['autoplay'] is None for v in all_videos)
        for v in all_videos:
            source   = "embedded" if v['embedded'] else "linked ⚠"
            ap_str   = ("? ⚠"   if v['autoplay'] is None
                        else "YES" if v['autoplay'] else "NO ⚠")
            loop_str = "?"   if v['loop']  is None else ("yes" if v['loop']  else "no")
            mute_str = "?"   if v['muted'] is None else ("yes" if v['muted'] else "no")
            sz_str   = fmt_size(v['size']) if v['size'] else "—"
            codec    = v.get('codec')
            if codec is None:
                codec_str = "—" if v['embedded'] else "n/a (linked)"
            elif codec in _CODEC_WARN:
                codec_str = f"{codec} ⚠"
            else:
                codec_str = codec
            print(f"  {v['slide']:>5}  {v['name']:<30}  {source:<9}  {codec_str:<14}"
                  f"  {ap_str:<8}  {loop_str:<5}  {mute_str:<5}  {sz_str}")
            if v.get('linked_path'):
                print(f"         ↳ {v['linked_path']}")
        no_autoplay_v = [v for v in all_videos if v['autoplay'] is False]
        if no_autoplay_v:
            print(f"\n  ⚠  {len(no_autoplay_v)} video(s) not set to autoplay — "
                  f"will require a manual click or trigger to start")
        if has_unknown:
            print(f"  ⚠  '?' = playback settings unreadable (older embed format) — "
                  f"verify autoplay/loop/mute manually in PowerPoint")
        unknown_codec = [v for v in all_videos if v['embedded'] and v.get('codec') is None]
        if unknown_codec and not _get_ffprobe():
            print(f"  ℹ  {len(unknown_codec)} codec(s) undetected — "
                  f"install FFmpeg (ffprobe) for full codec detection")

    # ── AUDIO ──────────────────────────────────────────────────────────────────
    def print_audio_section(items):
        section_header(f"AUDIO  ({len(items)})")
        if not items:
            print("  (none)")
            return
        for m in items:
            status   = "embedded" if m['embedded'] else "LINKED ⚠"
            sz       = f"  [{fmt_size(m['size'])}]" if m['size'] else ""
            fmt_warn = ("  ⚠ check codec support"
                        if m['ext'] not in SAFE_AUDIO_FORMATS else "")
            print(f"  slide {m['slide']:>2}  {status:<12}  "
                  f"{m['ext']:<6}  {m['name']}{sz}{fmt_warn}")
            if not m['embedded']:
                print(f"           ↳ linked path: {m['linked_path']}")

    print_audio_section(all_audios)

    # ── Images ─────────────────────────────────────────────────────────────────
    # Deduplicate by filename (same asset used on multiple slides)
    img_map = {}
    for img in all_images:
        k = img['name']
        if k not in img_map:
            img_map[k] = {**img, 'slides': [img['slide']]}
        else:
            img_map[k]['slides'].append(img['slide'])

    ext_counts     = defaultdict(int)
    total_img_size = 0
    for img in img_map.values():
        ext_counts[img['ext']] += 1
        if img['size']:
            total_img_size += img['size']

    linked_imgs = [img for img in img_map.values() if not img['embedded']]

    section_header(f"IMAGES  ({len(img_map)} unique)")
    for ext, count in sorted(ext_counts.items()):
        print(f"  {(ext or '(none)'):<10}  {count} file{'s' if count > 1 else ''}")
    if total_img_size:
        print(f"\n  Total image data   {fmt_size(total_img_size)}")
    if linked_imgs:
        print(f"\n  ⚠  {len(linked_imgs)} linked image(s) — must travel with the deck:")
        for img in linked_imgs:
            print(f"     slides {img['slides']}  {img['name']}")

    # ── Image resolution ───────────────────────────────────────────────────────
    ppi_issues, ppi_checked = image_ppi_issues(prs)

    section_header(f"IMAGE RESOLUTION  (target {_PPI_TARGET} PPI, range {_PPI_LOW}–{_PPI_HIGH} PPI)")
    if ppi_issues is None:
        print("  (install Pillow to enable: pip install Pillow)")
    elif ppi_checked == 0:
        print("  (no picture shapes found)")
    elif not ppi_issues:
        print(f"  ✓  {ppi_checked} image(s) checked — all within range")
    else:
        ok_count      = ppi_checked - len(ppi_issues)
        oversized     = [i for i in ppi_issues if i['too_high']]
        blurry        = [i for i in ppi_issues if i['too_low']]

        if verbose:
            # ── Full per-image table ──────────────────────────────────────────
            for issue in ppi_issues:
                dim_str  = f"{issue['img_w']}×{issue['img_h']} px"
                size_str = f"{issue['w_in']}\"×{issue['h_in']}\""
                if issue['too_high']:
                    pixel_ratio = round((issue['img_w'] * issue['img_h']) /
                                        (issue['ideal_w'] * issue['ideal_h']))
                    flag = (f"⚠ {issue['oversize_x']}× oversized "
                            f"(~{pixel_ratio}× more pixels than needed — "
                            f"reduce to ~{issue['ideal_w']}×{issue['ideal_h']} px, "
                            f"save ~{fmt_size(issue['file_size'] - issue['file_size'] // pixel_ratio)})")
                else:
                    flag = f"⚠ may appear blurry  (min {_PPI_LOW} PPI for standard displays)"
                print(f"  slide {issue['slide']:>3}  {issue['name']:<26}  "
                      f"{dim_str:<16}  {size_str:<12}  {issue['ppi']:>5} PPI  {flag}")
            print(f"\n  {ppi_checked} image(s) checked — "
                  f"{len(ppi_issues)} issue(s), {ok_count} ok")
        else:
            # ── Compact summary (default) ─────────────────────────────────────
            print(f"  {ppi_checked} image(s) checked — "
                  f"{len(ppi_issues)} issue(s)  ({ok_count} ok)")
            if oversized:
                over_slides = sorted({i['slide'] for i in oversized})
                worst = max(oversized, key=lambda i: i['oversize_x'] or 0)
                print(f"  ⚠  {len(oversized):>3} oversized  (>{_PPI_HIGH} PPI)  "
                      f"on slides {fmt_slide_list(over_slides, show_count=False)}")
                print(f"       worst: {worst['oversize_x']}× — {worst['name']}  "
                      f"slide {worst['slide']}  "
                      f"({worst['img_w']}×{worst['img_h']} px → "
                      f"~{worst['ideal_w']}×{worst['ideal_h']} px ideal)")
            if blurry:
                blur_slides = sorted({i['slide'] for i in blurry})
                print(f"  ⚠  {len(blurry):>3} blurry     (<{_PPI_LOW} PPI)   "
                      f"on slides {fmt_slide_list(blur_slides, show_count=False)}")
            print(f"\n  Add --verbose for per-image detail")

    # ── Transitions & animations ───────────────────────────────────────────────
    transitions  = []
    anim_slides  = []   # list of (slide_num, [dur_ms, ...])
    auto_advance = []

    for i, slide in enumerate(prs.slides, 1):
        t = slide_transition(slide, i)
        if t:
            transitions.append(t)
            if t['auto_ms']:
                auto_advance.append(i)
        durs = animation_effect_durations(slide.element)
        if durs:
            anim_slides.append((i, durs))

    master_anims = check_masters_for_animations(prs)

    section_header("TRANSITIONS & ANIMATIONS")
    real_transitions = [t for t in transitions if t['kind'] != 'none']
    print(f"  Transitions  ({len(real_transitions)}/{n_slides} slides)")
    used_default_dur = False
    for t in real_transitions:
        dur    = t['dur_ms']
        slow   = dur is not None and dur > TRANS_WARN_MS
        if dur is not None:
            d_str = fmt_ms(dur)
        else:
            d_str = "~0.7s *"
            used_default_dur = True
        s_tag  = f"  ⚠ SLOW (>{fmt_ms(TRANS_WARN_MS)})" if slow else ""
        adv    = "click" if t['click'] else "no-click"
        if t['auto_ms']:
            adv += f"  auto-advance {fmt_ms(int(t['auto_ms']))}"
        print(f"  slide {t['slide']:>3}  {t['kind']:<16}  {d_str:<12}{s_tag}  {adv}")
    if used_default_dur:
        print(f"  * duration not set — PowerPoint will use its built-in default (~0.7 s)")

    print(f"\n  Animations  ({len(anim_slides)}/{n_slides} slides)")
    for slide_n, durs in anim_slides:
        has_slow       = any(d is not None and d > ANIM_WARN_MS for d in durs)
        has_indefinite = any(d is None for d in durs)

        # Grouped summary: "19×2.0s ⚠  15×0.2s  1×1.0s"
        counts = Counter(durs)
        # Sort: None last, then by duration descending (longest / slowest first)
        sorted_keys = sorted(counts, key=lambda x: (x is None, -(x or 0)))
        dur_parts = []
        for d in sorted_keys:
            c = counts[d]
            if d is None:
                dur_parts.append(f"{c}×∞")
            else:
                flag = " ⚠" if d > ANIM_WARN_MS else ""
                dur_parts.append(f"{c}×{fmt_ms(d)}{flag}")

        flags = []
        if has_slow:
            flags.append(f"⚠ effect(s) > {fmt_ms(ANIM_WARN_MS)}")
        if has_indefinite:
            flags.append("⚠ indefinite duration")
        flag_str = "  " + "  ".join(flags) if flags else ""

        print(f"  slide {slide_n:>3}  {len(durs):>3} effect(s)   "
              f"{'  '.join(dur_parts)}{flag_str}")

    print(f"\n  Master / layouts")
    if master_anims:
        for label, durs in master_anims:
            print(f"  ⚠  {label}  has {len(durs)} animation(s) — "
                  f"will play on every slide using this master/layout")
    else:
        print(f"  ✓  No animations in slide master or layouts")

    if auto_advance:
        print(f"\n  ⚠  Auto-advancing slides: {fmt_slide_list(auto_advance)}  — verify timing on event system")

    # ── Speaker notes ──────────────────────────────────────────────────────────
    notes_slides = [
        i+1 for i, s in enumerate(prs.slides)
        if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip()
    ]
    section_header("SPEAKER NOTES")
    print(f"  {len(notes_slides)}/{n_slides} slides have notes" +
          (f"  {fmt_slide_list(notes_slides, show_count=False)}" if notes_slides else ""))

    # ── Hyperlinks ─────────────────────────────────────────────────────────────
    links = []
    for i, slide in enumerate(prs.slides, 1):
        for rel in slide.part.rels.values():
            if 'hyperlink' in rel.reltype.lower():
                links.append((i, rel.target_ref, rel.is_external))

    section_header(f"HYPERLINKS  ({len(links)})")
    if links:
        for slide_n, target, external in links:
            kind = "external" if external else "internal"
            print(f"  slide {slide_n:>2}  {kind:<10}  {target}")
    else:
        print("  (none)")

    # ── Embedded OLE objects ───────────────────────────────────────────────────
    ole_files = [f for f in zf.namelist() if '/embeddings/' in f]
    if ole_files:
        ole_total = sum(zf.getinfo(f).file_size for f in ole_files)
        section_header(f"EMBEDDED OBJECTS / OLE  ({len(ole_files)})")
        print(f"  {len(ole_files)} embedded object(s)   {fmt_size(ole_total)} total")
        print(f"  ⚠  Requires matching software (Excel, etc.) to render on playback machine")

    # ── Charts ─────────────────────────────────────────────────────────────────
    chart_files = [f for f in zf.namelist() if f.startswith('ppt/charts/chart')]
    if chart_files:
        section_header(f"CHARTS  ({len(chart_files)})")
        print(f"  {len(chart_files)} chart(s) — verify data labels render correctly on event system")

    # ── Preflight summary ──────────────────────────────────────────────────────
    issues   = []
    warnings = []

    linked_v = [v for v in all_videos if not v['embedded']]
    linked_a = [v for v in all_audios if not v['embedded']]

    def _slide_list(items, key='slide'):
        return fmt_slide_list(sorted({v[key] for v in items}), show_count=False)

    if linked_v:
        issues.append(
            f"{len(linked_v)} linked video(s) on slides {_slide_list(linked_v)} — "
            f"will break on a different machine"
        )
    if linked_a:
        issues.append(
            f"{len(linked_a)} linked audio file(s) on slides {_slide_list(linked_a)} — "
            f"will break on a different machine"
        )
    if linked_imgs:
        img_slides = sorted({s for img in linked_imgs for s in img['slides']})
        issues.append(
            f"{len(linked_imgs)} linked image(s) on slides {fmt_slide_list(img_slides, show_count=False)} — "
            f"will not display on a different machine"
        )

    # Codec warnings — group by codec so each gets its specific advisory
    _CODEC_FAIL = {'ProRes', 'WMV', 'WebM/MKV'}   # likely to fail outright
    _codec_groups = defaultdict(list)
    for v in all_videos:
        c = v.get('codec')
        if c and c in _CODEC_WARN:
            _codec_groups[c].append(v)
        elif v['embedded'] and v['ext'] not in SAFE_VIDEO_FORMATS and not c:
            _codec_groups[v['ext']].append(v)  # extension fallback
    for codec_label, vids in sorted(_codec_groups.items()):
        advice = _CODEC_WARN.get(codec_label, 'verify codec support on event system')
        msg = f"{len(vids)} {codec_label} video(s) on slides {_slide_list(vids)} — {advice}"
        if codec_label in _CODEC_FAIL:
            issues.append(msg)
        else:
            warnings.append(msg)

    no_autoplay_vids = [v for v in all_videos if v.get('autoplay') is False]
    if no_autoplay_vids:
        warnings.append(
            f"{len(no_autoplay_vids)} video(s) not set to autoplay on slides "
            f"{_slide_list(no_autoplay_vids)} — require manual click to start"
        )
    unknown_ap_vids = [v for v in all_videos if v.get('autoplay') is None]
    if unknown_ap_vids:
        warnings.append(
            f"{len(unknown_ap_vids)} video(s) with unreadable playback settings on slides "
            f"{_slide_list(unknown_ap_vids)} — verify autoplay/loop/mute manually in PowerPoint"
        )

    if not embedded_font_files and fonts:
        font_list = ', '.join(sorted(fonts)[:6])
        if len(fonts) > 6:
            font_list += f" (+{len(fonts)-6} more)"
        warnings.append(f"Fonts not embedded — must be installed on playback machine: {font_list}")

    if anim_slides:
        warnings.append(
            f"{len(anim_slides)} slide(s) have animations "
            f"{fmt_slide_list([s for s, _ in anim_slides], show_count=False)} — "
            f"test playback on event system"
        )

    slow_trans = [t for t in transitions
                  if t['dur_ms'] is not None and t['dur_ms'] > TRANS_WARN_MS]
    if slow_trans:
        details = ', '.join(
            f"slide {t['slide']} ({fmt_ms(t['dur_ms'])})" for t in slow_trans
        )
        warnings.append(f"Slow transition(s) > {fmt_ms(TRANS_WARN_MS)}: {details}")

    slow_anim_slides = [s for s, durs in anim_slides
                        if any(d is not None and d > ANIM_WARN_MS for d in durs)]
    if slow_anim_slides:
        warnings.append(
            f"Slow animation effect(s) > {fmt_ms(ANIM_WARN_MS)} on slides "
            f"{fmt_slide_list(slow_anim_slides)}"
        )

    indef_slides = [s for s, durs in anim_slides if any(d is None for d in durs)]
    if indef_slides:
        warnings.append(
            f"Indefinite-duration animation(s) on slides {fmt_slide_list(indef_slides)} — "
            f"may stall auto-advance"
        )

    if master_anims:
        labels = ', '.join(label for label, _ in master_anims)
        warnings.append(
            f"Animations in master/layout ({labels}) — "
            f"will play on every slide using that master/layout"
        )

    if auto_advance:
        warnings.append(
            f"Auto-advance timing on slides {fmt_slide_list(auto_advance)} — "
            f"verify on event system"
        )
    if ole_files:
        warnings.append(
            f"{len(ole_files)} OLE embedded object(s) — "
            f"requires matching software (Excel, etc.) to render"
        )
    if hidden_slides:
        warnings.append(
            f"{len(hidden_slides)} hidden slide(s): "
            f"{fmt_slide_list(hidden_slides, show_count=False)}"
        )
    if show_props['mode'] == 'kiosk (auto-advance, no menu)':
        warnings.append("Kiosk mode — presenter cannot skip slides manually")

    ext_links = [l for l in links if l[2]]
    if ext_links:
        ext_slides = fmt_slide_list(sorted({l[0] for l in ext_links}), show_count=False)
        warnings.append(
            f"{len(ext_links)} external hyperlink(s) on slides {ext_slides} — "
            f"confirm internet access on event system"
        )

    oversized_imgs = [i for i in (ppi_issues or []) if i['too_high']]
    blurry_imgs    = [i for i in (ppi_issues or []) if i['too_low']]
    if oversized_imgs:
        worst      = max(oversized_imgs, key=lambda i: i['oversize_x'] or 0)
        over_slides = fmt_slide_list(
            sorted({i['slide'] for i in oversized_imgs}), show_count=False
        )
        warnings.append(
            f"{len(oversized_imgs)} oversized image(s) on slides {over_slides} — "
            f"largest {worst['oversize_x']}× over (slide {worst['slide']}) — "
            f"downsample to reduce file size"
        )
    if blurry_imgs:
        warnings.append(
            f"{len(blurry_imgs)} blurry image(s) on slides "
            f"{fmt_slide_list(sorted({i['slide'] for i in blurry_imgs}), show_count=False)} — "
            f"may appear low-resolution on screen"
        )

    if not ratio_ok:
        warnings.append(
            f"Aspect ratio mismatch — deck is {aspect_label(w_emu, h_emu)}, "
            f"display is {disp_w}×{disp_h} ({aspect_label(disp_w, disp_h)}) — "
            f"will be {consequence}"
        )

    print(f"\n{'━'*60}")
    print("PREFLIGHT SUMMARY")
    print("─" * 60)
    if not issues and not warnings:
        print("  ✓  All checks passed — good to go")
    else:
        for item in issues:
            print(f"  ✗  FAIL   {item}")
        for item in warnings:
            print(f"  ⚠  WARN   {item}")
    print(f"{'━'*60}\n")

    zf.close()


if __name__ == '__main__':
    parser = argparse.ArgumentParser(
        description="AV preflight check for PowerPoint presentations.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument('file', help='Path to the .pptx file')
    parser.add_argument(
        '--display',
        metavar='WxH',
        type=parse_display,
        default=(1920, 1080),
        help='Target display resolution (default: 1920x1080)',
    )
    parser.add_argument(
        '--verbose', '-v',
        action='store_true',
        help='Show full per-image resolution detail (default: summary only)',
    )
    args = parser.parse_args()
    analyze(args.file, display_wh=args.display, verbose=args.verbose)
